#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
FRET timelapse → PowerPoint 자동 배열 (cm GUI 버전)

사용:
1) python FRET_timelapse_ppt_cm_gui.py 실행
2) 폴더 선택 (Sxx_txx_roiX_*.png 들어있는 crop 폴더)
3) "이미지 가로 (cm)" 입력 (세로는 비율 유지)
4) [PPT 생성] 클릭
→ 선택한 폴더 안에 FRET_timelapse_auto.pptx 생성

전제 파일명 예시:
  S03_t00_roi1_ch3.png        → BF/contrast
  S03_t00_roi1_DoverF.png     → FRET/ratio

FRET 판정 키워드: 'dov', 'ratio', 'fret'
BF   판정 키워드: 'ch', 'bf', 'phase', 'dic'
(필요시 classify_channel()에서 수정 가능)
"""

import os
import re
from collections import defaultdict

import tkinter as tk
from tkinter import filedialog, messagebox

from pptx import Presentation
from pptx.util import Inches

# ---------------- 공통 설정 ----------------

# 파일명 패턴: S03_t00_roi1_ch3.png / S03_t00_roi1_DoverF.png
FNAME_PATTERN = re.compile(
    r"^(S\d+)_t(\d+)_roi(\d+)_(.+)\.(png|tif|tiff)$",
    re.IGNORECASE
)

CM_PER_INCH = 2.54

def cm_to_inch(cm: float) -> float:
    return cm / CM_PER_INCH


def classify_channel(suffix: str):
    """채널명을 보고 FRET vs BF/Phase 구분."""
    suf = suffix.lower()
    # FRET / Ratio 계열
    if "dov" in suf or "ratio" in suf or "fret" in suf:
        return "fret"
    # BF / Phase / DIC / chX 계열
    if "bf" in suf or "phase" in suf or "dic" in suf or suf.startswith("ch"):
        return "bf"
    return None


def collect_pairs(img_dir):
    """
    img_dir에서 파일을 읽어
    {(stage, roi): [(time, fret_path, bf_path), ...]} 구조로 반환.
    """
    raw = defaultdict(dict)  # key: (stage, roi, time) -> {"fret": path, "bf": path}

    for fname in os.listdir(img_dir):
        fpath = os.path.join(img_dir, fname)
        if not os.path.isfile(fpath):
            continue

        m = FNAME_PATTERN.match(fname)
        if not m:
            continue

        stage, t_str, roi, suffix, _ext = m.groups()
        time_idx = int(t_str)
        ch_type = classify_channel(suffix)
        if ch_type is None:
            continue

        key = (stage, roi, time_idx)
        if key not in raw:
            raw[key] = {}
        raw[key][ch_type] = fpath

    timeline = defaultdict(list)

    # FRET + BF 페어 모두 있는 timepoint만 사용
    for (stage, roi, t), d in raw.items():
        if "fret" in d and "bf" in d:
            timeline[(stage, roi)].append((t, d["fret"], d["bf"]))

    # 시간 순 정렬
    for key in timeline:
        timeline[key].sort(key=lambda x: x[0])

    return timeline


def build_ppt(timeline, img_dir, img_width_cm):
    """
    timeline을 기반으로 PPT 생성.

    - img_width_cm: 사용자가 지정한 '이미지 한 장 가로(cm)'
    - 슬라이드 크기: 16:9 고정 (약 33.867cm x 19.05cm)
    - 필요한 경우, 모든 이미지를 동일 비율로 자동 축소하여 슬라이드 폭에 맞춤
    """
    if not timeline:
        return False, "유효한 FRET/BF 페어가 없습니다.\n파일명 패턴과 키워드를 확인하세요."

    prs = Presentation()

    # 슬라이드 크기 16:9로 고정 (원하는 경우 여기만 바꾸면 됨)
    slide_w_in = cm_to_inch(33.867)  # 약 13.33 inch
    slide_h_in = cm_to_inch(19.05)   # 약 7.5 inch
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)

    blank = prs.slide_layouts[6]  # 빈 슬라이드

    # 여백/간격 (cm 기준)
    left_margin_cm = 1.0
    top_margin_cm = 1.5
    row_gap_cm = 0.3    # FRET와 BF 사이 세로 간격
    col_gap_cm = 0.1    # timepoint 사이 가로 간격

    left_margin = Inches(cm_to_inch(left_margin_cm))
    top_margin = Inches(cm_to_inch(top_margin_cm))
    row_gap = Inches(cm_to_inch(row_gap_cm))
    col_gap = Inches(cm_to_inch(col_gap_cm))

    slide_width = prs.slide_width

    desired_img_w = Inches(cm_to_inch(img_width_cm))

    # Stage, ROI 정렬용
    def sort_key(item):
        (stage, roi) = item[0]
        stage_num = int(stage[1:])  # 'S03' → 3
        roi_num = int(roi)
        return (stage_num, roi_num)

    for (stage, roi), seq in sorted(timeline.items(), key=sort_key):
        if not seq:
            continue

        slide = prs.slides.add_slide(blank)
        n = len(seq)

        # 우선 사용자가 지정한 폭으로 배치했을 때 필요한 총 너비
        total_gap = col_gap * (n - 1) if n > 1 else 0
        needed_width = left_margin * 2 + desired_img_w * n + total_gap

        # 슬라이드보다 크면 비율 유지해서 전체 축소
        if needed_width > slide_width:
            scale = (slide_width - left_margin * 2 - total_gap) / (desired_img_w * n)
            if scale <= 0:
                return False, (
                    f"{stage} ROI{roi}: 이미지 수가 너무 많아서 배치 불가.\n"
                    "이미지 가로(cm)를 줄이거나 timepoint를 줄여주세요."
                )
            img_w = desired_img_w * scale
        else:
            img_w = desired_img_w

        # 세로 위치 (정사각형 가정: 가로 길이 기준)
        fret_top = top_margin
        bf_top = fret_top + img_w + row_gap

        for idx, (t, fret_path, bf_path) in enumerate(seq):
            left = left_margin + idx * (img_w + col_gap)

            # 위: FRET / ratio
            slide.shapes.add_picture(fret_path, left, fret_top, width=img_w)
            # 아래: BF / contrast
            slide.shapes.add_picture(bf_path, left, bf_top, width=img_w)

        # 텍스트 라벨
        label_box = slide.shapes.add_textbox(
            Inches(cm_to_inch(1.0)), Inches(cm_to_inch(0.5)),
            Inches(cm_to_inch(15)), Inches(cm_to_inch(1.0))
        )
        tf = label_box.text_frame
        tf.text = f"{stage}  ROI{roi}  (위: FRET / 아래: BF, t00 → t{seq[-1][0]:02d})"

    out_path = os.path.join(img_dir, "FRET_timelapse_auto.pptx")
    prs.save(out_path)
    return True, f"FRET_timelapse_auto.pptx 생성 완료\n폴더: {out_path}"


# ---------------- GUI 부분 ----------------

def browse_folder():
    path = filedialog.askdirectory(title="FRET / BF PNG 폴더 선택")
    if path:
        folder_var.set(path)

def run_generator():
    img_dir = folder_var.get().strip()
    if not img_dir:
        messagebox.showwarning("확인", "이미지 폴더를 선택하세요.")
        return
    if not os.path.isdir(img_dir):
        messagebox.showerror("오류", "유효한 폴더가 아닙니다.")
        return

    try:
        img_w_cm = float(img_width_cm_var.get())
    except ValueError:
        messagebox.showwarning("확인", "이미지 가로(cm)를 숫자로 입력하세요.")
        return

    if img_w_cm <= 0:
        messagebox.showwarning("확인", "이미지 가로(cm)는 0보다 커야 합니다.")
        return

    timeline = collect_pairs(img_dir)
    success, msg = build_ppt(timeline, img_dir, img_w_cm)

    if success:
        messagebox.showinfo("완료", msg)
    else:
        messagebox.showerror("실패", msg)


if __name__ == "__main__":
    root = tk.Tk()
    root.title("FRET Timelapse PPT 자동 생성기 (cm)")

    # 폴더 선택
    tk.Label(root, text="이미지 폴더:").grid(row=0, column=0, sticky="e", padx=5, pady=5)
    folder_var = tk.StringVar()
    tk.Entry(root, textvariable=folder_var, width=50).grid(row=0, column=1, padx=5, pady=5)
    tk.Button(root, text="찾아보기", command=browse_folder).grid(row=0, column=2, padx=5, pady=5)

    # 이미지 가로(cm)
    tk.Label(root, text="이미지 1개 가로 (cm):").grid(row=1, column=0, sticky="e", padx=5, pady=5)
    img_width_cm_var = tk.StringVar(value="2.0")  # 기본값: 2 cm (원하는대로 변경)
    tk.Entry(root, textvariable=img_width_cm_var, width=10).grid(
        row=1, column=1, sticky="w", padx=5, pady=5
    )

    # 실행 버튼
    tk.Button(root, text="PPT 생성", command=run_generator, width=25).grid(
        row=2, column=0, columnspan=3, pady=15
    )

    root.mainloop()
